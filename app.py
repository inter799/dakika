# -*- coding: utf-8 -*-
"""
文件上传 → 题目生成器（纯本地模式，无需登录、无需API）
- 内置纯本地规则生成题目（无需任何 API Key）
- 支持 txt / doc / docx / pdf / rtf / md / html / json / csv / xlsx / pptx / jpg / png / gif 等文件
"""

import os, json, re, base64, zipfile, shutil
from io import BytesIO
from pathlib import Path
from datetime import datetime

from flask import Flask, render_template, request, jsonify
from flask_cors import CORS

# ==================== 可选依赖检测 ====================
try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from PIL import Image
    HAS_PIL = True
except ImportError:
    HAS_PIL = False

try:
    import pytesseract
    HAS_TESSERACT = True
except ImportError:
    HAS_TESSERACT = False

# PDF 解析（可选）
try:
    import PyPDF2
    HAS_PYPDF2 = True
except ImportError:
    HAS_PYPDF2 = False

try:
    import pdfplumber
    HAS_PDFPLUMBER = True
except ImportError:
    HAS_PDFPLUMBER = False

# Excel / PPT 解析（可选）
try:
    import openpyxl
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# ==================== 平台检测 ====================
import platform as _platform
_IS_WINDOWS = _platform.system() == 'Windows'
_IS_SERVERLESS = bool(os.environ.get('RAILWAY_ENVIRONMENT') or os.environ.get('NETLIFY') or os.environ.get('RENDER'))

# ==================== 初始化 ====================
app = Flask(__name__)
app.config['JSON_AS_ASCII'] = False  # 关键：确保中文JSON输出不被转义为\uXXXX
CORS(app)

# 纯本地模式：所有请求共用一个数据目录（无需登录、无用户隔离）
if _IS_SERVERLESS:
    DATA_ROOT = Path('/tmp/data')
    UPLOAD_FOLDER = Path('/tmp/uploads')
    BANK_FILE = Path('/tmp/question_bank.json')
else:
    DATA_ROOT = Path(os.environ.get('DATA_DIR', Path(__file__).parent / 'data'))
    UPLOAD_FOLDER = DATA_ROOT / 'uploads'
    BANK_FILE = DATA_ROOT / 'question_bank.json'

DATA_ROOT.mkdir(parents=True, exist_ok=True)
UPLOAD_FOLDER.mkdir(parents=True, exist_ok=True)


# ==================== 题库存储辅助（纯本地共享） ====================
def load_bank():
    if BANK_FILE.exists():
        return json.loads(BANK_FILE.read_text(encoding='utf-8'))
    return _new_bank()


def _new_bank():
    return {"spaces": {}, "questions": [], "next_id": 1}


def save_bank(bank):
    BANK_FILE.write_text(json.dumps(bank, ensure_ascii=False, indent=2), encoding='utf-8')


def _generate_space_id(name):
    """根据文档名生成唯一 space_id"""
    base = re.sub(r'[^a-zA-Z0-9\u4e00-\u9fa5]', '_', name)[:30]
    ts = datetime.now().strftime('%Y%m%d%H%M%S')
    return f"{base}_{ts}"


ALLOWED_EXTENSIONS = {
    # 文本文档
    'txt', 'md', 'json', 'csv', 'xml', 'log', 'yaml', 'yml',
    # 富文本文档
    'rtf', 'html', 'htm',
    # Word
    'doc', 'docx',
    # PDF
    'pdf',
    # Excel
    'xlsx', 'xlsm',
    # PowerPoint
    'pptx',
    # 图片
    'jpg', 'jpeg', 'png', 'bmp', 'gif'
}
MAX_FILE_SIZE = 10 * 1024 * 1024  # 10MB


def safe_filename(filename):
    """保留中文等非ASCII字符的安全文件名处理"""
    unsafe_chars = '<>:"/\\|?*'
    name = filename
    for ch in unsafe_chars:
        name = name.replace(ch, '_')
    name = name.strip(' .')
    if len(name) > 200:
        base, ext = (name.rsplit('.', 1) if '.' in name else (name, ''))
        name = base[:195] + ('.' + ext if ext else '')
    return name or 'unnamed_file'


# ==================== 文本提取 ====================
def extract_text_from_txt(filepath):
    for enc in ['utf-8', 'utf-16', 'utf-16le', 'utf-16be', 'gbk', 'gb2312', 'latin-1']:
        try:
            with open(filepath, 'r', encoding=enc) as f:
                return f.read()
        except UnicodeDecodeError:
            continue
    try:
        with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
            return f.read()
    except Exception:
        return "[错误] 无法解码该文本文件"


def _is_zip_docx(filepath):
    """检测文件是否为标准的 zip 格式 docx"""
    try:
        with zipfile.ZipFile(filepath, 'r') as z:
            return 'word/document.xml' in z.namelist()
    except Exception:
        return False


def extract_text_from_docx(filepath):
    if not HAS_DOCX:
        return "[错误] 请安装 python-docx: pip install python-docx"

    if not _is_zip_docx(filepath):
        return "[错误] 该文件不是标准的 .docx 格式。可能是旧版 .doc 文件改名，请用 Word 打开后另存为 .docx，或转换为 .txt 上传。"

    try:
        doc = DocxDocument(filepath)
        paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        paras.append(cell.text.strip())
        text = '\n'.join(paras)
        if not text.strip():
            return "[提示] 该 docx 文件内容为空，或文字位于图片/文本框中无法提取。"
        return text
    except Exception as e:
        return f"[错误] 解析docx失败: {e}"


def extract_text_from_doc(filepath):
    """.doc 旧版二进制格式：尝试提取可读文本"""
    com_errors = []
    try:
        if _is_zip_docx(filepath):
            return extract_text_from_docx(filepath)

        try:
            with open(filepath, 'rb') as f:
                raw = f.read()

            text_parts = []
            for enc in ['utf-8', 'gbk', 'gb2312', 'utf-16-le', 'latin-1']:
                try:
                    decoded = raw.decode(enc, errors='ignore')
                    paragraphs = re.findall('[\\u4e00-\\u9fa5a-zA-Z0-9\\s，。！？；：、""''《》（）…—.,!?;:()\\[\\]【】]{20,}', decoded)
                    if paragraphs:
                        text_parts.extend(paragraphs)
                        break
                except Exception:
                    continue

            if text_parts:
                result = '\n\n'.join(text_parts)
                if len(result) >= 100:
                    return result
        except Exception:
            pass

        word = None
        try:
            import pythoncom
            pythoncom.CoInitialize()
        except Exception as e:
            com_errors.append(f"COM初始化失败: {e}")

        try:
            import win32com.client as win32
            for progid in ["Word.Application", "KWPS.Application", "WPS.Application", "Kwps.Application", "Et.Application"]:
                try:
                    word = win32.Dispatch(progid)
                    break
                except Exception as e:
                    com_errors.append(f"{progid}: {e}")
                    word = None

            if word:
                word.Visible = False
                word.DisplayAlerts = 0
                abs_path = os.path.abspath(filepath)
                doc = word.Documents.Open(abs_path, ReadOnly=True)
                text = doc.Content.Text
                doc.Close(SaveChanges=False)
                extracted = text.strip() if text else ''
                try:
                    word.Quit()
                except Exception as e:
                    com_errors.append(f"Word.Application.Quit: {e}")
                if extracted:
                    return extracted
        except Exception as e:
            com_errors.append(f"COM提取失败: {e}")
        finally:
            try:
                if word:
                    word.Quit()
            except Exception:
                pass
            try:
                import pythoncom
                pythoncom.CoUninitialize()
            except Exception:
                pass

        import subprocess
        try:
            proc = subprocess.run(['antiword', str(filepath)], capture_output=True, text=True, timeout=10)
            if proc.returncode == 0 and proc.stdout.strip():
                return proc.stdout.strip()
        except (FileNotFoundError, subprocess.TimeoutExpired):
            pass

        err_detail = '; '.join(com_errors) if com_errors else '未知原因'
        return (f"[错误] 无法解析旧版 .doc 文件。\n"
                f"[建议] 请用 Word/WPS 打开后另存为 .docx 或 .txt 再上传。\n"
                f"[详情] {err_detail}")
    except Exception as e:
        return f"[错误] 无法解析旧版 .doc 文件: {e}。请确保电脑已安装 Word 或 WPS，或者手动另存为 .docx / .txt 后重新上传。"


def extract_text_from_image(filepath):
    if not HAS_PIL:
        return "[错误] 请安装 Pillow: pip install Pillow"
    try:
        img = Image.open(filepath)
        if img.mode != 'RGB':
            img = img.convert('RGB')
        if HAS_TESSERACT:
            text = pytesseract.image_to_string(img, lang='chi_sim+eng')
            return text if text.strip() else "[提示] 图片中未检测到文字"
        else:
            buffered = BytesIO()
            img.save(buffered, format='PNG')
            return "[IMAGE_BASE64]" + base64.b64encode(buffered.getvalue()).decode('utf-8')
    except Exception as e:
        return f"[错误] 解析图片失败: {e}"


def extract_text_from_pdf(filepath):
    """解析 PDF 文件，优先使用 pdfplumber，回退到 PyPDF2"""
    if HAS_PDFPLUMBER:
        try:
            pages_text = []
            with pdfplumber.open(filepath) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        pages_text.append(page_text)
            text = '\n\n'.join(pages_text)
            if text.strip():
                return text
        except Exception:
            pass

    if HAS_PYPDF2:
        try:
            pages_text = []
            with open(filepath, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        pages_text.append(page_text)
            text = '\n\n'.join(pages_text)
            if text.strip():
                return text
            return "[提示] PDF 中未检测到可提取文字（可能是扫描版图片PDF，需要 OCR）"
        except Exception as e:
            return f"[错误] 解析 PDF 失败: {e}。请安装依赖: py -m pip install pdfplumber PyPDF2"

    return "[错误] 请安装 PDF 解析依赖: py -m pip install pdfplumber PyPDF2"


def extract_text_from_rtf(filepath):
    """解析 RTF 富文本格式：剥离控制符，保留可见文本"""
    try:
        with open(filepath, 'rb') as f:
            raw = f.read()
        for enc in ['utf-8', 'gbk', 'gb2312', 'latin-1']:
            try:
                text = raw.decode(enc, errors='ignore')
                break
            except Exception:
                text = ''
        if not text:
            return "[错误] 无法读取 RTF 文件"
        text = re.sub(r'\\[a-z]+\d*\s?', '', text)
        text = re.sub(r'\\[*\{\}\~\-]', '', text)
        text = re.sub(r'\{\\[^}]+\}', '', text)
        text = re.sub(r'[\{\}]', '', text)
        text = re.sub(r'\\par\b', '\n', text, flags=re.IGNORECASE)
        text = re.sub(r'\\tab\b', '\t', text, flags=re.IGNORECASE)
        text = re.sub(r'\\line\b', '\n', text, flags=re.IGNORECASE)
        text = re.sub(r'\\[0-9a-fA-F]{2}', '', text)
        text = re.sub(r'\\$', '', text)
        text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', text)
        text = '\n'.join(line.strip() for line in text.splitlines() if line.strip())
        return text if text.strip() else "[提示] RTF 文件中未检测到文字"
    except Exception as e:
        return f"[错误] 解析 RTF 失败: {e}"


def extract_text_from_plain(filepath):
    """读取普通文本类文件：md / html / htm / json / csv / xml / log / yaml / yml"""
    try:
        with open(filepath, 'rb') as f:
            raw = f.read()
        for enc in ['utf-8', 'gbk', 'gb2312', 'latin-1']:
            try:
                text = raw.decode(enc)
                break
            except UnicodeDecodeError:
                text = raw.decode('utf-8', errors='replace')
        return text
    except Exception as e:
        return f"[错误] 读取文本文件失败: {e}"


def extract_text_from_xlsx(filepath):
    """解析 Excel 表格文件，读取每个单元格的文本"""
    if not HAS_OPENPYXL:
        return "[错误] 请安装 openpyxl: py -m pip install openpyxl"
    try:
        wb = openpyxl.load_workbook(filepath, data_only=True)
        lines = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = ' '.join(str(cell) for cell in row if cell is not None)
                if row_text.strip():
                    lines.append(row_text)
        text = '\n'.join(lines)
        return text if text.strip() else "[提示] Excel 文件中未检测到文字"
    except Exception as e:
        return f"[错误] 解析 Excel 失败: {e}"


def extract_text_from_pptx(filepath):
    """解析 PowerPoint 演示文稿"""
    if not HAS_PPTX:
        return "[错误] 请安装 python-pptx: py -m pip install python-pptx"
    try:
        prs = Presentation(filepath)
        texts = []
        for slide_idx, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    texts.append(shape.text.strip())
        text = '\n\n'.join(texts)
        return text if text.strip() else "[提示] PPT 中未检测到文字"
    except Exception as e:
        return f"[错误] 解析 PPT 失败: {e}"


def extract_text(filepath, original_filename):
    ext = original_filename.rsplit('.', 1)[1].lower() if '.' in original_filename else ''
    handlers = {
        'docx': extract_text_from_docx,
        'doc': extract_text_from_doc,
        'pdf': extract_text_from_pdf,
        'rtf': extract_text_from_rtf,
        'xlsx': extract_text_from_xlsx,
        'xlsm': extract_text_from_xlsx,
        'pptx': extract_text_from_pptx,
        'jpg': extract_text_from_image, 'jpeg': extract_text_from_image,
        'png': extract_text_from_image, 'bmp': extract_text_from_image,
        'gif': extract_text_from_image,
        'txt': extract_text_from_plain,
        'md': extract_text_from_plain,
        'json': extract_text_from_plain,
        'csv': extract_text_from_plain,
        'xml': extract_text_from_plain,
        'log': extract_text_from_plain,
        'yaml': extract_text_from_plain,
        'yml': extract_text_from_plain,
        'html': extract_text_from_plain,
        'htm': extract_text_from_plain,
    }
    return handlers.get(ext, lambda _: f"[错误] 不支持 .{ext} 格式")(filepath)


# ==================== 文本规范化 ====================
def normalize_text(text):
    if not text:
        return text
    text = text.translate(str.maketrans(
        'ＡＢＣＤＥＦＧＨＩＪＫＬＭＮＯＰＱＲＳＴＵＶＷＸＹＺ'
        'ａｂｃｄｅｆｇｈｉｊｋｌｍｎｏｐｑｒｓｔｕｖｗｘｙｚ'
        '０１２３４５６７８９',
        'ABCDEFGHIJKLMNOPQRSTUVWXYZ'
        'abcdefghijklmnopqrstuvwxyz'
        '0123456789'
    ))
    text = text.replace('．', '.').replace('：', ':').replace('，', ',')
    text = text.replace('（', '(').replace('）', ')').replace('【', '[').replace('】', ']')
    text = text.replace('《', '<').replace('》', '>').replace('；', ';')
    text = text.replace('？', '?').replace('！', '!').replace('～', '~')
    text = text.replace('\u3000', ' ')
    text = text.replace('\r\n', '\n').replace('\r', '\n')
    text = re.sub(r'\n{3,}', '\n\n', text)
    lines = [line.strip() for line in text.split('\n')]
    text = '\n'.join(lines)
    return text


# ==================== 题目识别与解析 ====================
_Q_SPLIT = re.compile(
    r'(?:^|\n)\s*'
    r'(?:'
    r'\d+[\.、．\)）]\s*'
    r'|[(（]\s*\d+\s*[)）]\s*'
    r'|第\s*\d+\s*题\s*'
    r'|第\s*[一二三四五六七八九十]+\s*题\s*'
    r')',
    re.MULTILINE
)

_OPT_RE = re.compile(r'([A-E])[\.、．\)）]\s*')
_OPT_PAREN_RE = re.compile(r'[（(]\s*([A-E])\s*[)）]\s*')
_STEM_INLINE_ANSWER_RE = re.compile(
    r'[（(]\s*([A-Ea-e]+)\s*[)）]|'
    r'【\s*答案\s*[：:]\s*([A-Ea-e]+)\s*】|'
    r'[（(]\s*(正确|错误|对|错)\s*[)）]',
    re.IGNORECASE
)

_ANS_RE = re.compile(
    r'(?:标准答案|正确答案|参考答案|答案|建议答案|正确选项)'
    r'[是为]?\s*[：:]\s*[（(]?\s*'
    r'([A-Ea-e]+(?:[、，,\s/]+[A-Ea-e]+)*)\s*[)）]?'
    r'|(?:标准答案|正确答案|参考答案|答案|建议答案)[是为]?\s*[：:]\s*(正确|错误|对|错)'
    r'|[【\[]\s*(?:标准答案|正确答案|参考答案|答案)\s*[】\]]\s*[（(]?\s*([A-Ea-e]+)\s*[)）]?',
    re.IGNORECASE
)
_ANS_TF_RE = re.compile(r'(?:答案|正确答案)[：:是为]?\s*(正确|错误|对|错)')
_ANS_LETTER_RE = re.compile(r'(?:答案|正确答案|参考答案)[：:是为]?\s*[（(]?\s*([A-Ea-e]+)\s*[)）]?')


def _count_question_numbers(text):
    return len(_Q_SPLIT.findall(text))


def detect_existing_questions(text):
    if not text or len(text) < 20:
        return False, 0, ''
    q_count = _count_question_numbers(text)
    if q_count < 1:
        return False, 0, ''
    opt_count = len(_OPT_RE.findall(text)) + len(_OPT_PAREN_RE.findall(text))
    ans_count = len(_ANS_RE.findall(text)) + len(_ANS_TF_RE.findall(text))
    if q_count >= 1 and (opt_count >= 2 or ans_count >= 1):
        if opt_count >= ans_count:
            return True, q_count, 'choice'
        return True, q_count, 'mixed'
    return False, 0, ''


def parse_questions_from_text(text, question_type='choice'):
    questions = []
    parts = _Q_SPLIT.split(text)
    if not parts or len(parts) < 2:
        q = _parse_one_block(text.strip(), 1)
        return [q] if q else []

    q_num = 0
    for i in range(1, len(parts)):
        block = parts[i].strip()
        if not block or len(block) < 4:
            continue
        q_num += 1
        q = _parse_one_block(block, q_num)
        if q:
            questions.append(q)

    if not questions:
        questions = _parse_tf_questions(text)

    return questions


def _parse_one_block(block, q_id):
    if len(block) < 4:
        return None

    answer = ''
    ans_pos = len(block)

    inline_m = _STEM_INLINE_ANSWER_RE.search(block)
    if inline_m:
        raw = inline_m.group(1) or inline_m.group(2) or inline_m.group(3)
        raw = (raw or '').strip()
        text_before = block[:inline_m.start()]
        if re.search(r'[A-E][\.、．\)）]', text_before[-200:] if len(text_before) > 200 else text_before):
            answer = ''
        else:
            if raw in ('正确', '对'):
                answer = '正确'
            elif raw in ('错误', '错'):
                answer = '错误'
            else:
                letters = sorted(set(re.findall(r'[A-Ea-e]', raw)))
                answer = ''.join(letters).upper()
            block = block[:inline_m.start()] + block[inline_m.end():]

    if not answer:
        m_ans = _ANS_RE.search(block)
        if m_ans:
            raw = m_ans.group(1) or m_ans.group(3) or ''
            tf = m_ans.group(2) or ''
            if tf:
                answer = '正确' if tf in ('正确', '对') else '错误'
            elif raw:
                answer = ''.join(re.findall(r'[A-Ea-e]', raw)).upper()
                if len(answer) > 1:
                    answer = ''.join(sorted(set(answer)))
            if answer:
                ans_pos = m_ans.start()

    if not answer:
        m_tf = _ANS_TF_RE.search(block)
        if m_tf:
            tf = m_tf.group(1)
            answer = '正确' if tf in ('正确', '对') else '错误'
            ans_pos = m_tf.start()

    options = []
    opt_start = len(block)

    for m in _OPT_RE.finditer(block):
        if m.start() >= ans_pos:
            break
        letter = m.group(1).upper()
        opt_start = min(opt_start, m.start())
        content_start = m.end()
        next_m = _OPT_RE.search(block, content_start)
        content_end = next_m.start() if (next_m and next_m.start() < ans_pos) else ans_pos
        content = block[content_start:content_end].strip()
        options.append((letter, content))

    if not options:
        for m in _OPT_PAREN_RE.finditer(block):
            if m.start() >= ans_pos:
                break
            letter = m.group(1).upper()
            opt_start = min(opt_start, m.start())
            content_start = m.end()
            next_m = _OPT_PAREN_RE.search(block, content_start)
            content_end = next_m.start() if (next_m and next_m.start() < ans_pos) else ans_pos
            content = block[content_start:content_end].strip()
            options.append((letter, content))

    if options:
        stem = block[:opt_start]
    else:
        stem = block[:ans_pos] if ans_pos < len(block) else block

    stem = _clean_stem(stem)

    if not stem or len(stem.strip()) < 1:
        return None

    qtype = _detect_type(block, answer, options)

    opt_list = [f'{letter}. {content}' for letter, content in options]

    if not opt_list and qtype == 'tf':
        opt_list = ['正确. 正确', '错误. 错误']
    elif not opt_list and qtype in ('choice', 'multi'):
        opt_list = ['A. 选项A', 'B. 选项B', 'C. 选项C', 'D. 选项D']

    answer_guessed = False
    if not answer:
        answer_guessed = True
        if qtype == 'tf':
            answer = '正确' if q_id % 2 == 0 else '错误'
        else:
            answer = options[0][0] if options else 'A'

    analysis = _gen_analysis(qtype, stem, answer, opt_list)

    result = {
        'id': q_id,
        'type': qtype,
        'stem': stem.strip(),
        'options': opt_list,
        'answer': answer,
        'analysis': analysis,
    }
    if answer_guessed:
        result['answer_guessed'] = True
    return result


def _clean_stem(stem):
    stem = _ANS_RE.sub('', stem)
    stem = _ANS_TF_RE.sub('', stem)
    stem = _ANS_LETTER_RE.sub('', stem)
    stem = _STEM_INLINE_ANSWER_RE.sub('', stem)
    stem = _OPT_RE.sub(' ', stem)
    stem = _OPT_PAREN_RE.sub(' ', stem)
    stem = re.sub(r'\s+[（(]?\s*[A-Ea-e]\s*[)）]?\s*$', '', stem)
    stem = re.sub(r'(?:标准答案|正确答案|参考答案|答案|建议答案)\s*[：:是为]?\s*$', '', stem)
    stem = re.sub(r'\s+', ' ', stem)
    stem = stem.strip().rstrip('.,，。、;；:：')
    return stem


def _detect_type(block, answer, options):
    if len(options) == 2:
        texts = ' '.join(o[1] for o in options if isinstance(o, tuple) and len(o) > 1)
        if re.search(r'正确|错误|对|错', texts):
            if answer in ('正确', '错误', '对', '错'):
                return 'tf'

    if len(answer) >= 2 and re.match(r'^[A-E]+$', answer):
        return 'multi'

    if re.search(r'多选|不定项|多项|至少.*选|选出.?包括', block):
        return 'multi'

    if len(options) >= 5 and len(answer) >= 2 and re.match(r'^[A-E]+$', answer):
        return 'multi'

    return 'choice'


def _gen_analysis(qtype, stem, answer, opt_list):
    if not opt_list:
        return '请参考原文内容。'

    ans_letters = list(answer) if answer else []
    correct_texts = []
    for opt in opt_list:
        parts = opt.split('.', 1) if '.' in opt else (opt[:1], opt[1:])
        letter = parts[0].strip()
        content = parts[1].strip() if len(parts) > 1 else opt
        if letter in ans_letters:
            correct_texts.append(content)

    correct_str = '；'.join(correct_texts[:3]) if correct_texts else '相关选项'

    if qtype == 'tf':
        state = '符合' if answer == '正确' else '不符合'
        return f'该陈述{state}原文内容。'
    elif qtype == 'multi':
        return f'本题为多选题，共 {len(ans_letters)} 个正确选项。正确答案 {answer}：{correct_str}。'
    else:
        return f'正确答案为 {answer}：{correct_str}。'


def _parse_tf_questions(text):
    questions = []

    blocks = re.split(
        r'(?:^|\n)\s*(?:\d+[\.、．\)]\s*|[(（]\d+[)）]\s*|第[一二三四五六七八九十\d]+题\s*)',
        text
    )
    if len(blocks) < 2:
        return []

    for idx, block in enumerate(blocks[1:], start=1):
        block = block.strip()
        if not block or len(block) < 8:
            continue

        answer = ""
        ans_patterns = [
            r'(?:答案|正确答案)[：:是为]?\s*(正确|错误|对|错|True|False|T|F)',
            r'[（(]\s*(正确|错误|对|错|[×✓✔✗✘√])\s*[)）]',
            r'(正确|错误|对|错)\s*[（(]\s*[)）]',
        ]
        for pat in ans_patterns:
            m = re.search(pat, block)
            if m:
                raw = m.group(1)
                if raw in ('正确', '对', '✓', '✔', '√', 'True', 'T', 'true', 't'):
                    answer = "正确"
                elif raw in ('错误', '错', '×', '✗', '✘', 'False', 'F', 'false', 'f'):
                    answer = "错误"
                break

        stem = block
        for pat in ans_patterns:
            stem = re.sub(pat, '', stem).strip()
        stem = re.sub(r'\s*[（(]\s*[)）]\s*', '', stem).strip()

        if not stem or len(stem) < 4:
            continue

        analysis = ""
        am = re.search(r'(?:解析|分析)[：:]\s*(.+?)(?=\n|$)', block)
        if am:
            analysis = am.group(1).strip()

        questions.append({
            "id": len(questions) + 1,
            "type": "tf",
            "stem": stem,
            "answer": answer if answer else ("正确" if idx % 2 == 0 else "错误"),
            "analysis": analysis if analysis else "从原文解析提取"
        })

    return questions


def _guess_answer(options):
    return options[0][0] if options else "A"


# ==================== 智能题数估算 ====================
def estimate_question_count(text, max_count=50):
    text_len = len(text)
    sentences = re.split(r'[。！？\n；;]+', text)
    sentence_count = len([s for s in sentences if len(s.strip()) >= 10])
    estimated = max(3, min(max_count, sentence_count // 2))
    return estimated


# ==================== 纯本地模式（无需API） ====================
def _split_sentences(text):
    sents = re.split(r'[。！？\n；;]+', text)
    return [s.strip() for s in sents if len(s.strip()) >= 12]


def _extract_keywords(text, topn=5):
    cn_words = re.findall(r'[\u4e00-\u9fa5]{2,4}', text)
    en_words = re.findall(r'[a-zA-Z]{3,}', text)
    all_words = cn_words + en_words

    freq = {}
    for w in all_words:
        if w.lower() not in {'这是', '不是', '一个', '可以', '进行', '使用', '没有',
                              '他们', '我们', '这个', '那个', '什么', '因为', '所以',
                              '而且', '但是', '如果', '就是', '已经', '还是', '或者'}:
            freq[w] = freq.get(w, 0) + 1

    sorted_words = sorted(freq.items(), key=lambda x: x[1], reverse=True)
    return [w for w, _ in sorted_words[:topn]]


def _extract_numbers(text):
    patterns = [
        r'(\d+\.?\d*\s*[%％])',
        r'(\d{4}\s*年)',
        r'([约近超达]?\s*\d+\.?\d*\s*[万亿千百]?\s*[个只条项次元美元吨米克])'
    ]
    results = []
    for p in patterns:
        results.extend(re.findall(p, text))
    return results[:5]


def generate_choice_local(text, count):
    sents = _split_sentences(text)
    if not sents:
        return []

    keywords = _extract_keywords(text, 10)
    numbers = _extract_numbers(text)

    max_q = min(len(sents), 50) if count == 0 else min(count, len(sents))

    questions = []
    for i in range(max_q):
        sent = sents[i % len(sents)]
        words = re.findall(r'[\u4e00-\u9fa5a-zA-Z]{2,}', sent)
        if not words:
            words = re.findall(r'\S+', sent)
        if not words:
            words = ['核心要素']

        kw = words[i % len(words)]

        templates = [
            {
                "stem": f'根据原文内容，"{sent[:45]}..."，以下哪项说法是正确的？',
                "options": [
                    f"A. {kw}是核心要素",
                    f"B. 相关内容与{kw}无关",
                    f"C. 以上说法均不正确",
                    f"D. 需要结合上下文综合判断"
                ],
                "answer": "A",
                "analysis": f'原文"{sent[:80]}..."明确指出{kw}的关键地位。'
            },
            {
                "stem": f'"{sent[:45]}..."，在这段内容中主要强调的是？',
                "options": [
                    f"A. 方法论的重要性",
                    f"B. {kw}的核心作用",
                    f"C. 外部因素的影响",
                    f"D. 历史背景的分析"
                ],
                "answer": "B",
                "analysis": f'内容重点围绕"{kw}"展开论述。'
            },
            {
                "stem": f'以下哪项最符合原文中"{sent[:35]}..."的表述？',
                "options": [
                    f"A. 原文明确否定此观点",
                    f"B. {kw}在其中起决定性作用",
                    f"C. 没有足够信息判断",
                    f"D. 此观点需要更多证据支持"
                ],
                "answer": "B",
                "analysis": f'根据上下文，"{kw}"是核心概念。'
            },
            {
                "stem": f'关于"{sent[:45]}..."，以下理解最准确的是？',
                "options": [
                    f"A. 该描述主要指{kw}",
                    f"B. 该描述缺乏明确指向",
                    f"C. 该描述与主题无关",
                    f"D. 该描述存在歧义"
                ],
                "answer": "A",
                "analysis": f'综合分析，原文讨论的核心是{kw}。'
            }
        ]

        tmpl = templates[i % len(templates)]
        questions.append({
            "id": len(questions) + 1,
            "stem": tmpl["stem"],
            "options": tmpl["options"],
            "answer": tmpl["answer"],
            "analysis": tmpl["analysis"]
        })

        if len(questions) >= count:
            break

    return questions


def generate_essay_local(text, count):
    sents = _split_sentences(text)
    if not sents:
        return []

    keywords = _extract_keywords(text, 15)
    questions = []

    for i in range(count):
        sent = sents[i % len(sents)]
        kw = keywords[i % len(keywords)] if keywords else "相关内容"
        kw2 = keywords[(i + 1) % len(keywords)] if keywords else "相关概念"

        templates = [
            {
                "stem": f'请简述"{sent[:60]}..."的主要内容，并说明{kw}在其中的作用。',
                "answer": sent if len(sent) < 300 else sent[:300] + "...",
                "key_points": [
                    f"准确概括核心内容",
                    f"阐述{kw}的定义与内涵",
                    f"分析{kw}与整体内容的关系",
                    f"给出自己的理解或总结"
                ],
                "difficulty": "medium"
            },
            {
                "stem": f'结合原文分析，{kw}与{kw2}之间存在怎样的关系？请举例说明。',
                "answer": f"根据原文内容，{kw}和{kw2}存在相互关联。{sent[:200]}",
                "key_points": [
                    f"明确{kw}的定义",
                    f"分析{kw2}的内涵",
                    f"阐述两者之间的关系",
                    f"举出原文中的具体例子"
                ],
                "difficulty": "hard"
            },
            {
                "stem": f'阅读"{sent[:50]}..."，请用简洁的语言概括其表达的核心观点（不超过100字）。',
                "answer": f"核心观点：{sent[:150]}",
                "key_points": [
                    "准确提取核心观点",
                    "表述清晰简洁",
                    "不超过字数限制",
                    "不遗漏关键信息"
                ],
                "difficulty": "easy"
            }
        ]

        tmpl = templates[i % len(templates)]
        questions.append({
            "id": len(questions) + 1,
            "stem": tmpl["stem"],
            "answer": tmpl["answer"],
            "key_points": tmpl["key_points"],
            "difficulty": tmpl["difficulty"]
        })

        if len(questions) >= count:
            break

    return questions


def generate_fill_local(text, count):
    sents = _split_sentences(text)
    keywords = _extract_keywords(text, 20)

    questions = []
    for i in range(count):
        sent = sents[i % len(sents)]
        kw = keywords[i % len(keywords)] if keywords else "___"
        blank_sent = sent.replace(kw, "___", 1) if kw in sent else sent[:40] + "___" + sent[40:80]
        questions.append({
            "id": len(questions) + 1,
            "stem": blank_sent[:150],
            "answer": kw,
            "analysis": f'根据原文语境，此处应填入"{kw}"。'
        })
        if len(questions) >= count:
            break
    return questions


def generate_tf_local(text, count):
    sents = _split_sentences(text)
    if not sents:
        return []

    keywords = _extract_keywords(text, 10)

    max_q = min(len(sents), 50) if count == 0 else min(count, len(sents))

    questions = []
    for i in range(max_q):
        sent = sents[i % len(sents)]
        kw = keywords[i % len(keywords)] if keywords else "核心要素"
        is_true = i % 2 == 0

        if is_true:
            stem = f'{kw}{"" if kw in sent else "在原文内容中"}可以找到相关依据。'
            answer = "正确"
            analysis = f'原文"{sent[:60]}..."证实了这一点。'
        else:
            opposite_kw = keywords[(i + 1) % len(keywords)] if len(keywords) > 1 else "无关因素"
            stem = f'{opposite_kw}是影响{kw}的唯一决定因素。'
            answer = "错误"
            analysis = f'原文表明{kw}受多种因素影响，并非仅由{opposite_kw}决定。'

        questions.append({
            "id": len(questions) + 1,
            "stem": stem,
            "answer": answer,
            "analysis": analysis
        })
    return questions


def generate_questions_local(text, question_type, question_count=5):
    generators = {
        "choice": generate_choice_local,
        "tf": generate_tf_local,
        "fill": generate_fill_local,
        "essay": generate_essay_local,
    }

    if question_type == 'all':
        if question_count == 0:
            per_type = max(3, min(15, estimate_question_count(text, max_count=50) // 4))
        else:
            per_type = max(2, question_count // 4)
        grouped = {}
        total = 0
        for t in ['choice', 'tf', 'fill', 'essay']:
            arr = generators[t](text, per_type)
            grouped[t] = arr
            total += len(arr)
        return {
            "error": None,
            "grouped": grouped,
            "total": total,
            "provider_used": "本地模式（无需API）",
        }

    gen = generators.get(question_type, generate_choice_local)
    questions = gen(text, question_count)
    return {
        "error": None,
        "questions": questions,
        "provider_used": "本地模式（无需API）",
    }


# ==================== 路由 ====================
@app.after_request
def add_utf8_charset(response):
    content_type = response.headers.get('Content-Type', '')
    if content_type and 'charset' not in content_type:
        if 'text/html' in content_type:
            response.headers['Content-Type'] = 'text/html; charset=utf-8'
        elif 'application/json' in content_type:
            response.headers['Content-Type'] = 'application/json; charset=utf-8'
        elif 'text/plain' in content_type:
            response.headers['Content-Type'] = 'text/plain; charset=utf-8'
    return response


@app.route('/')
def index():
    return render_template('index.html')


@app.route('/api/upload', methods=['POST'])
def upload_file():
    """统一的文件上传接口，支持 multipart/form-data 和 JSON+base64 两种方式"""
    try:
        filename = None
        filepath = None

        if 'file' in request.files:
            file = request.files['file']
            if not file.filename:
                return jsonify({"error": "文件名为空"}), 400
            ext = file.filename.rsplit('.', 1)[1].lower() if '.' in file.filename else ''
            if ext not in ALLOWED_EXTENSIONS:
                return jsonify({"error": f"不支持 .{ext}，支持: {', '.join(ALLOWED_EXTENSIONS)}"}), 400
            filename = safe_filename(file.filename)
            filepath = UPLOAD_FOLDER / filename
            file.seek(0, os.SEEK_END)
            if file.tell() > MAX_FILE_SIZE:
                return jsonify({"error": "文件超过 10MB"}), 400
            file.seek(0)
            file.save(str(filepath))

        elif request.is_json:
            data = request.get_json(silent=True)
            if data and data.get('file_base64') and data.get('filename'):
                filename = safe_filename(data['filename'])
                ext = filename.rsplit('.', 1)[1].lower() if '.' in filename else ''
                if ext not in ALLOWED_EXTENSIONS:
                    return jsonify({"error": f"不支持 .{ext}"}), 400
                try:
                    file_bytes = base64.b64decode(data['file_base64'])
                except Exception:
                    return jsonify({"error": "Base64 解码失败"}), 400
                if len(file_bytes) > MAX_FILE_SIZE:
                    return jsonify({"error": "文件超过 10MB"}), 400
                filepath = UPLOAD_FOLDER / filename
                with open(str(filepath), 'wb') as f:
                    f.write(file_bytes)
            else:
                return jsonify({"error": "JSON 数据缺少 file_base64 或 filename"}), 400
        else:
            return jsonify({"error": "没有上传文件"}), 400

        if not filepath or not filename:
            return jsonify({"error": "文件保存失败"}), 500

        full_text = extract_text(str(filepath), filename)
        if full_text.startswith("[错误]"):
            return jsonify({"error": full_text.strip("[]错误 ").strip()}), 400
        if full_text.startswith("[提示]"):
            return jsonify({"error": full_text.strip("[]提示 ").strip()}), 400
        full_text = normalize_text(full_text)
        truncated = len(full_text) > 6000
        display_text = full_text[:6000] + ("\n\n...[文本过长已截断]..." if truncated else "")

        has_questions, detected_count, detected_type = detect_existing_questions(full_text)
        parsed_questions = None
        if has_questions:
            parsed = parse_questions_from_text(full_text, detected_type)
            if parsed:
                parsed_questions = {
                    "type": detected_type,
                    "questions": parsed,
                    "count": len(parsed)
                }

        return jsonify({
            "filename": filename,
            "text": display_text,
            "full_text": full_text,
            "full_length": len(full_text),
            "truncated": truncated,
            "has_existing_questions": has_questions,
            "parsed_questions": parsed_questions
        })
    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({"error": f"服务器处理上传失败: {str(e)}"}), 500


@app.route('/api/generate', methods=['POST'])
def generate_questions():
    try:
        data = request.get_json()
        if not data:
            return jsonify({"error": "请求体不是有效的 JSON"}), 400

        text = data.get('text', '')
        question_type = data.get('type', 'choice')
        question_count = int(data.get('count', 5))

        valid_types = ('choice', 'tf', 'fill', 'essay', 'all')
        if question_type not in valid_types:
            question_type = 'all'

        if not text.strip():
            return jsonify({"error": "文本内容为空"}), 400

        text = normalize_text(text)

        has_existing, _, detected_type = detect_existing_questions(text)
        if has_existing:
            parsed = parse_questions_from_text(text, detected_type)
            if parsed:
                return jsonify({
                    "error": None,
                    "questions": parsed,
                    "total": len(parsed),
                    "provider_used": "原文档解析（保持原答案）",
                    "notice": "已检测到文件中包含现成题目，直接解析并保留原答案。"
                })

        is_all_mode = (question_count == 0)
        if is_all_mode:
            question_count = estimate_question_count(text, max_count=50)
        else:
            question_count = min(question_count, 50)

        result = generate_questions_local(text, question_type, question_count)
        return jsonify(result)
    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({"error": f"生成题目失败: {str(e)}"}), 500


# ==================== 题库管理 API（文档空间） ====================
@app.route('/api/bank', methods=['GET'])
def get_bank():
    bank = load_bank()
    spaces_summary = []
    for sid, sdata in bank.get("spaces", {}).items():
        wrong_book = sdata.get("wrong_book", {})
        spaces_summary.append({
            "space_id": sid,
            "name": sdata.get("name", sid),
            "created_at": sdata.get("created_at", ""),
            "question_count": len(sdata.get("questions", [])),
            "question_types": list(set(q.get("type", "choice") for q in sdata.get("questions", []))),
            "wrong_count": len(wrong_book),
        })
    spaces_summary.sort(key=lambda x: x.get("created_at", ""), reverse=True)
    return jsonify({
        "spaces": spaces_summary,
        "total_spaces": len(spaces_summary),
        "total_questions": sum(s.get("question_count", 0) for s in spaces_summary),
        "questions": bank.get("questions", []),
    })


@app.route('/api/bank/spaces', methods=['GET'])
def list_spaces():
    bank = load_bank()
    result = []
    for sid, sdata in bank.get("spaces", {}).items():
        result.append({
            "space_id": sid,
            "name": sdata.get("name", sid),
            "created_at": sdata.get("created_at", ""),
            "question_count": len(sdata.get("questions", [])),
        })
    result.sort(key=lambda x: x.get("created_at", ""), reverse=True)
    return jsonify(result)


@app.route('/api/bank/space/<space_id>', methods=['GET'])
def get_space(space_id):
    bank = load_bank()
    space = bank.get("spaces", {}).get(space_id)
    if not space:
        return jsonify({"error": "空间未找到"}), 404
    return jsonify({
        "space_id": space_id,
        "name": space.get("name", space_id),
        "created_at": space.get("created_at", ""),
        "questions": space.get("questions", []),
        "total": len(space.get("questions", []))
    })


@app.route('/api/bank/save-batch', methods=['POST'])
def save_batch_to_bank():
    data = request.get_json()
    if not data:
        return jsonify({"error": "请求数据为空"}), 400

    questions = data.get('questions', [])
    space_name = data.get('space_name', '未命名文档')
    space_id = data.get('space_id', None)
    source_text = data.get('source_text', '')

    if not questions or not isinstance(questions, list):
        return jsonify({"error": "没有题目数据"}), 400

    bank = load_bank()

    if not space_id:
        space_id = _generate_space_id(space_name)

    existing = bank.get("spaces", {}).get(space_id, None)
    if existing:
        start_id = existing.get("next_id", len(existing.get("questions", [])) + 1)
        for q in questions:
            q["bank_id"] = start_id
            start_id += 1
            existing["questions"].append(q)
        existing["next_id"] = start_id
        if source_text:
            existing["source_text"] = source_text
    else:
        start_id = 1
        for q in questions:
            q["bank_id"] = start_id
            start_id += 1
        bank["spaces"][space_id] = {
            "name": space_name,
            "created_at": datetime.now().isoformat(),
            "questions": list(questions),
            "next_id": start_id,
            "source_text": source_text
        }

    save_bank(bank)
    return jsonify({
        "success": True,
        "space_id": space_id,
        "space_name": space_name,
        "added": len(questions),
        "total_in_space": len(bank["spaces"][space_id]["questions"])
    })


@app.route('/api/bank', methods=['POST'])
def add_to_bank():
    data = request.get_json()
    if not data:
        return jsonify({"error": "请求数据为空"}), 400

    questions = data.get('questions', [])
    if not isinstance(questions, list):
        questions = [data]
    if not questions:
        return jsonify({"error": "没有题目数据"}), 400

    space_id = data.get('space_id') or data.get('space_name', '未命名文档')
    space_name = data.get('space_name', '未命名文档')
    source_text = data.get('source_text', '')

    if not re.search(r'\d{14}$', space_id or ''):
        space_id = None

    bank = load_bank()

    if not space_id:
        space_id = _generate_space_id(space_name)

    existing = bank.get("spaces", {}).get(space_id)
    if existing:
        start_id = existing.get("next_id", len(existing.get("questions", [])) + 1)
        for q in questions:
            q["bank_id"] = start_id
            start_id += 1
            existing["questions"].append(q)
        existing["next_id"] = start_id
        if source_text:
            existing["source_text"] = source_text
    else:
        start_id = 1
        for q in questions:
            q["bank_id"] = start_id
            start_id += 1
        bank["spaces"][space_id] = {
            "name": space_name,
            "created_at": datetime.now().isoformat(),
            "questions": list(questions),
            "next_id": start_id,
            "source_text": source_text
        }

    save_bank(bank)
    total = len(bank["spaces"][space_id]["questions"])
    return jsonify({
        "success": True,
        "space_id": space_id,
        "space_name": space_name,
        "added": len(questions),
        "total": total,
        "total_in_space": total
    })


@app.route('/api/bank/space/<space_id>', methods=['PUT'])
def rename_space(space_id):
    data = request.get_json()
    if not data:
        return jsonify({"error": "请求数据为空"}), 400

    new_name = data.get('name', '').strip()
    if not new_name:
        return jsonify({"error": "空间名称不能为空"}), 400

    bank = load_bank()
    space = bank.get("spaces", {}).get(space_id)
    if not space:
        return jsonify({"error": "空间未找到"}), 404
    space["name"] = new_name
    save_bank(bank)
    return jsonify({"success": True, "space_id": space_id, "name": new_name})


@app.route('/api/bank/space/<space_id>', methods=['DELETE'])
def delete_space(space_id):
    bank = load_bank()
    if space_id not in bank.get("spaces", {}):
        return jsonify({"error": "空间未找到"}), 404
    del bank["spaces"][space_id]
    save_bank(bank)
    return jsonify({"success": True})


# ==================== 空间错题本 API ====================
@app.route('/api/bank/space/<space_id>/wrong', methods=['POST'])
def record_wrong_answers(space_id):
    bank = load_bank()
    space = bank.get("spaces", {}).get(space_id)
    if not space:
        return jsonify({"error": "空间未找到"}), 404

    data = request.get_json()
    wrong_list = data.get('wrong_questions', [])
    if not wrong_list:
        return jsonify({"error": "没有错题数据"}), 400

    if "wrong_book" not in space:
        space["wrong_book"] = {}

    for wq in wrong_list:
        bank_id = str(wq.get("bank_id", ""))
        if bank_id:
            space["wrong_book"][bank_id] = {
                "bank_id": wq["bank_id"],
                "stem": wq.get("stem", ""),
                "type": wq.get("type", "choice"),
                "options": wq.get("options", []),
                "answer": wq.get("answer", ""),
                "analysis": wq.get("analysis", ""),
                "user_answer": wq.get("user_answer", ""),
                "recorded_at": datetime.now().isoformat()
            }

    save_bank(bank)
    return jsonify({
        "success": True,
        "wrong_count": len(space["wrong_book"]),
        "space_id": space_id
    })


@app.route('/api/bank/space/<space_id>/wrong', methods=['GET'])
def get_wrong_questions(space_id):
    bank = load_bank()
    space = bank.get("spaces", {}).get(space_id)
    if not space:
        return jsonify({"error": "空间未找到"}), 404

    wrong_book = space.get("wrong_book", {})
    wrong_list = list(wrong_book.values())
    wrong_list.sort(key=lambda x: x.get("recorded_at", ""), reverse=True)

    return jsonify({
        "space_id": space_id,
        "space_name": space.get("name", space_id),
        "wrong_questions": wrong_list,
        "wrong_count": len(wrong_list)
    })


@app.route('/api/bank/space/<space_id>/wrong/clear', methods=['POST'])
def clear_wrong_questions(space_id):
    bank = load_bank()
    space = bank.get("spaces", {}).get(space_id)
    if not space:
        return jsonify({"error": "空间未找到"}), 404

    space["wrong_book"] = {}
    save_bank(bank)
    return jsonify({"success": True, "space_id": space_id})


@app.route('/api/bank/<int:bank_id>', methods=['DELETE'])
def delete_from_bank(bank_id):
    bank = load_bank()
    for sid, sdata in bank.get("spaces", {}).items():
        before = len(sdata["questions"])
        sdata["questions"] = [q for q in sdata["questions"] if q.get("bank_id") != bank_id]
        if len(sdata["questions"]) < before:
            if not sdata["questions"]:
                del bank["spaces"][sid]
            save_bank(bank)
            return jsonify({"success": True, "deleted_from": sid})
    return jsonify({"error": "题目未找到"}), 404


@app.route('/api/bank/clear', methods=['POST'])
def clear_bank():
    save_bank(_new_bank())
    return jsonify({"success": True, "total": 0})


@app.route('/api/config', methods=['GET'])
def get_config():
    return jsonify({
        "has_docx": HAS_DOCX,
        "has_pil": HAS_PIL,
        "has_tesseract": HAS_TESSERACT,
        "has_pdfplumber": HAS_PDFPLUMBER,
        "has_pypdf2": HAS_PYPDF2,
        "has_openpyxl": HAS_OPENPYXL,
        "has_pptx": HAS_PPTX,
        "allowed_types": list(ALLOWED_EXTENSIONS),
        "max_file_size_mb": MAX_FILE_SIZE // (1024 * 1024),
    })


if __name__ == '__main__':
    _OS_INFO = f"Windows {_platform.release()}" if _IS_WINDOWS else _platform.platform()
    print("=" * 60)
    print("  文件上传 → 题目生成器（纯本地模式）")
    print(f"  运行平台: {_OS_INFO}")
    print("  访问: http://127.0.0.1:5000")
    print("  无需登录，无需API，打开即用")
    print("=" * 60)
    # Railway/Render 等平台通过 gunicorn 启动时会读取 $PORT
    port = int(os.environ.get('PORT', 5000))
    app.run(debug=False, host='0.0.0.0', port=port)
